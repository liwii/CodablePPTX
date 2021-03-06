import yaml
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from copy import deepcopy

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def get_theme(theme, key, index):
    if key not in theme:
        return {}
    if len(theme[key]) <= index:
        return {}
    return theme[key][index]


def add_shape(base, content, theme):
    default_config = {'type': 'rectangle'}
    config = {**theme, **content}
    if config['type'] == 'rectangle':
        shape_type = MSO_SHAPE.RECTANGLE
    elif config['type'] == 'rounded_rectangle':
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        raise ValueError(f'Invalide shape type value: {config["type"]}')
    l, t, w, h = frame_to_position(config['frame'])
    shapes = base.shapes
    shape = shapes.add_shape(shape_type, l, t, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(config['color'][1:])
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.shadow.visible = False


def add_image(base, content, theme):
    default_config = {'aspect': 'fill'}
    if isinstance(content, dict):
        content_config = content
    else:
        content_config = {'file': content}
    config = {**default_config, **theme, **content_config}
    img = Image.open(config['file'])
    l, t, w, h = frame_to_position(config['frame'])
    frame_ratio = h / w
    img_w, img_h = img.size
    img_ratio = img_h / img_w
    if config['aspect'] == 'fill':
        left = l
        top = t
        height = h
        width = w
        if frame_ratio > img_ratio:
            crop_horizontal = 0.5 - (img_ratio / frame_ratio) / 2
            crop_vertical = 0
        else:
            crop_horizontal = 0
            crop_vertical = 0.5 - (frame_ratio / img_ratio) / 2
    elif config['aspect'] == 'fit':
        crop_horizontal = 0
        crop_vertical = 0
        if frame_ratio > img_ratio:
            left = l
            width = w
            height = width * img_ratio
            top = t + (h - height) / 2
        else:
            top = t
            height = h
            width = height / img_ratio
            left = l + (w - width) / 2
    else:
        raise ValueError(f'Invalid aspect value: {config["aspect"]}')
    picture = base.shapes.add_picture(config['file'],
                                      left,
                                      top,
                                      width=width,
                                      height=height)
    picture.crop_right = crop_horizontal
    picture.crop_left = crop_horizontal
    picture.crop_top = crop_vertical
    picture.crop_bottom = crop_vertical


def frame_to_position(frame):
    xl, xh = [float(f) for f in frame['x']]
    yl, yh = [float(f) for f in frame['y']]

    left = xl * SLIDE_WIDTH / 100
    width = xh * SLIDE_WIDTH / 100 - left
    top = yl * SLIDE_HEIGHT / 100
    height = yh * SLIDE_HEIGHT / 100 - top
    return left, top, width, height


def add_text(base, content, theme):
    default_config = {
        'value': '',
        'font': 'Arial',
        'fontsize': 26,
        'fontcolor': "#4D4D4D",
        'valign': 'top',
        'halign': 'left',
    }
    if isinstance(content, dict):
        content_config = content
    else:
        content_config = {'value': content}
    config = {**default_config, **theme, **content_config}
    l, t, w, h = frame_to_position(config['frame'])
    tx_box = base.shapes.add_textbox(l, t, w, h)
    tf = tx_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config['value']
    p.font.name = config['font']
    p.font.size = Pt(int(config['fontsize']))
    p.font.color.rgb = RGBColor.from_string(config['fontcolor'][1:])

    if config['valign'] == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif config['valign'] == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif config['valign'] == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        raise ValueError(f'Invalid valign value: {config["valign"]}')

    if config['halign'] == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif config['halign'] == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif config['halign'] == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        raise ValueError(f'Invalid halign value: {config["halign"]}')


def ignore_fixed(content, theme):
    for k in theme:
        if isinstance(theme[k], str) or isinstance(theme[k], dict):
            continue
        filled_in, not_filled_in = [], []
        for box in theme[k]:
            if ('fixed' in box) and (box['fixed'] != 'false'):
                filled_in.append(box)
            else:
                not_filled_in.append(box)
        theme[k] = filled_in + not_filled_in
        if k in content:
            content[k] = ([{}] * len(filled_in)) + content[k]


def add_empty_config(content, theme, key):
    if key not in content:
        content[key] = []
    if key not in theme:
        theme[key] = []
    diff = len(content[key]) - len(theme[key])
    if diff > 0:
        longer, shorter = content[key], theme[key]
    else:
        longer, shorter = theme[key], content[key]
    for i in range(abs(diff)):
        shorter.append({})


def nth_block(segment, margin, n, num_blocks):
    seg_begin, seg_end = int(segment[0]), int(segment[1])
    segment_len = seg_end - seg_begin
    block_len = (segment_len - margin * (num_blocks - 1)) / num_blocks
    new_seg_begin = n * (block_len + margin) + seg_begin
    return [new_seg_begin, new_seg_begin + block_len]


def fill_hstack_content(stack_contents, stack_frame, content, theme):
    config = {**stack_contents, **stack_frame}
    key = config['type']
    blocks = len(stack_contents[key])
    for i, stack_content in enumerate(stack_contents[key]):
        frame = deepcopy(config['template'])
        frame['frame'] = {}
        frame['frame']['y'] = config['frame']['y']
        frame['frame']['x'] = nth_block(config['frame']['x'],
                                        int(config['margin']), i, blocks)
        if key == 'vstack':
            fill_vstack_content(stack_content, frame, content, theme)
        elif key == 'hstack':
            fill_hstack_content(stack_content, frame, content, theme)
        else:
            content[key].append(stack_content)
            theme[key].append(frame)


def fill_vstack_content(stack_contents, stack_frame, content, theme):
    config = {**stack_contents, **stack_frame}
    key = config['type']
    blocks = len(stack_contents[key])
    for i, stack_content in enumerate(stack_contents[key]):
        frame = deepcopy(config['template'])
        frame['frame'] = {}
        frame['frame']['x'] = config['frame']['x']
        frame['frame']['y'] = nth_block(config['frame']['y'],
                                        int(config['margin']), i, blocks)
        if key == 'vstack':
            fill_vstack_content(stack_content, frame, content, theme)
        elif key == 'hstack':
            fill_hstack_content(stack_content, frame, content, theme)
        else:
            content[key].append(stack_content)
            theme[key].append(frame)


def fill_zip(content, theme, key):
    elements_content = content[key]
    elements_theme = theme[key]
    if len(elements_content) > len(elements_theme):
        elements_theme += [{}] * len(elements_content)
    else:
        elements_content += [{}] * len(elements_theme)
    return zip(elements_content, elements_theme)


def apply_design(base, content, theme, pageidx):
    ignore_fixed(content, theme)
    for key in [
            'title', 'subtitle', 'text', 'image', 'shape', 'pagenum', 'vstack',
            'hstack'
    ]:
        add_empty_config(content, theme, key)

    for stack_contents, stack_frame in zip(content['vstack'], theme['vstack']):
        fill_vstack_content(stack_contents, stack_frame, content, theme)

    for stack_contents, stack_frame in zip(content['hstack'], theme['hstack']):
        #breakpoint()
        fill_hstack_content(stack_contents, stack_frame, content, theme)

    for shape, shape_theme in fill_zip(content, theme, 'shape'):
        add_shape(base, shape, shape_theme)

    for title, title_theme in fill_zip(content, theme, 'title'):
        add_text(base, title, title_theme)

    for subtitle, subtitle_theme in fill_zip(content, theme, 'subtitle'):
        add_text(base, subtitle, subtitle_theme)

    for text, text_theme in fill_zip(content, theme, 'text'):
        add_text(base, text, text_theme)

    for image, image_theme in fill_zip(content, theme, 'image'):
        add_image(base, image, image_theme)

    for pagenum, pagenum_theme in fill_zip(content, theme, 'pagenum'):
        pagenum['value'] = str(pageidx + 1)
        add_text(base, pagenum, pagenum_theme)


def apply_master(layout, master):
    for master_content in master:
        if master_content == 'id':
            continue
        if master_content in layout:
            layout[master_content] = master[master_content] + layout[
                master_content]
        else:
            layout[master_content] = master[master_content]
    return layout


def draw_theme(theme, layouts):
    for layout in layouts:
        if layout['id'] == theme:
            return deepcopy(layout)
    raise ValueError(f'Theme doesn\'t exist: {theme}')


def draw_master(master_title, masters):
    for master in masters:
        if master['id'] == master_title:
            return master
    raise ValueError(f'Master doesn\'t exist: {master_title}')


def main(filename, output):
    f = open(filename, 'r')
    presentation_dict = yaml.load(f, Loader=yaml.BaseLoader)
    f.close()
    slides = presentation_dict['slides']
    layouts = []
    if 'import' not in presentation_dict:
        presentation_dict['import'] = []

    for layout_file in presentation_dict['import']:
        layout_f = open(layout_file, 'r')
        layout_dict = yaml.load(layout_f, Loader=yaml.BaseLoader)
        layout_f.close()
        layouts_template = layout_dict['layouts']
        if 'masters' in layout_dict:
            masters = layout_dict['masters']
            for i, layout in enumerate(layouts_template):
                if 'master' in layout:
                    master = draw_master(layout['master'], masters)
                    layouts_template[i] = apply_master(layout, master)
        layouts = layouts + layouts_template
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]
    for i, slide in enumerate(slides):
        blank_slide = prs.slides.add_slide(blank_slide_layout)
        if 'theme' in slide:
            theme = draw_theme(slide['theme'], layouts)
        else:
            theme = {}
        apply_design(blank_slide, slide, theme, i)
    prs.save(output)


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
